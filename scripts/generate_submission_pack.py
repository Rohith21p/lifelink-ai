#!/usr/bin/env python3
"""
Generate complete college project submission pack for LifeLink AI.
Outputs:
- Final Project Documentation (.docx + .pdf)
- Final Presentation (.pptx + .pdf)
- Project Synopsis (.docx + .pdf)
- One-page Summary (.pdf)
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

PROJECT_TITLE = "LifeLink AI - Intelligent Donor-Patient Matching and Healthcare Coordination System"
TODAY = date.today().strftime("%d %B %Y")


@dataclass
class Chapter:
    title: str
    paragraphs: list[str]
    bullets: list[str]


CHAPTERS: list[Chapter] = [
    Chapter(
        "1. Introduction",
        [
            "Healthcare systems often operate in silos where patient records, donor information, and blood inventory are stored in separate tools. This separation increases response time for critical cases and limits the ability of hospital coordinators to make evidence-based decisions quickly.",
            "LifeLink AI is designed as an integrated healthcare SaaS platform that centralizes patient management, donor management, match review workflows, report metadata, notifications, and blood bank operations. The platform focuses on operational intelligence, traceability, and reliable data exchange across modules.",
            "The project follows a full-stack architecture using Next.js, NestJS, PostgreSQL, and Docker. The approach balances practical usability for hospital teams with an extensible backend that can incorporate future AI decision support without rewriting core workflows.",
        ],
        [
            "Unified coordination platform for donor-patient workflows",
            "Rule-based matching foundation with explainable scoring",
            "Operational dashboard for executive and coordinator roles",
            "Modular architecture for future AI enhancement",
        ],
    ),
    Chapter(
        "2. Problem Statement",
        [
            "Hospitals handling organ and blood coordination face delays due to fragmented data and manual follow-up processes. Coordinators frequently rely on spreadsheets or disconnected tools to compare donor compatibility, track urgency, and communicate updates.",
            "In emergency scenarios, delayed matching and inconsistent communication can directly affect patient outcomes. Existing workflows often lack transparent score reasoning, centralized notification logs, and a live view of blood stock levels.",
            "The core problem is the absence of a single, intelligent workflow system that allows hospitals to register, evaluate, and monitor donor-patient coordination activities in real time.",
        ],
        [
            "Fragmented patient and donor data",
            "Slow manual compatibility checks",
            "Limited visibility into low blood stock alerts",
            "No standardized end-to-end review audit trail",
        ],
    ),
    Chapter(
        "3. Existing System",
        [
            "Conventional hospital environments use EMR systems primarily for clinical records, while coordination tasks are managed through calls, messaging, and manual trackers. These systems are not optimized for donor-patient workflow automation.",
            "Most existing tools provide either patient management or inventory management, but not both in a unified workflow with compatibility scoring and actionable review states.",
            "As a result, operational teams spend significant effort on data reconciliation instead of case decision quality and timeline management.",
        ],
        [
            "Heavy dependence on manual communication channels",
            "No integrated compatibility score workflow",
            "Weak auditability for review actions and notifications",
            "Minimal dashboard-driven decision support",
        ],
    ),
    Chapter(
        "4. Proposed System",
        [
            "LifeLink AI proposes a modular full-stack system that connects patient registration, donor registration, matching review, notifications, reports, and blood bank inventory in one platform.",
            "The solution introduces rule-based compatibility scoring as a foundation layer. Coordinators can review score breakdowns, add notes, move match states through a defined lifecycle, and trigger communication events through mocked adapters.",
            "The system architecture is designed for incremental improvement. Existing modules are production-style and database-backed, while AI extraction and communication adapters remain replaceable interfaces for future deployment-grade integrations.",
        ],
        [
            "End-to-end coordinator workflow",
            "Database-first design with strict schema relations",
            "API-first backend with Swagger and DTO validation",
            "Frontend UX optimized for healthcare operations",
        ],
    ),
    Chapter(
        "5. Objectives",
        [
            "The primary objective is to build an intelligent coordination platform that improves match review speed, data visibility, and case tracking reliability.",
            "A secondary objective is to establish a scalable technical base that supports advanced AI modules, legal workflow automation, and production integrations in future phases.",
        ],
        [
            "Reduce donor-patient matching turnaround time",
            "Provide explainable compatibility scores",
            "Enable transparent activity logs and notifications",
            "Improve blood inventory visibility and low-stock response",
            "Deliver a robust full-stack codebase with clean architecture",
        ],
    ),
    Chapter(
        "6. Scope",
        [
            "This project scope covers donor-patient coordination workflows inside a hospital network demo context. It includes patient and donor CRUD operations, rule-based match scoring, review workflows, report metadata handling, notification abstractions, and blood bank inventory monitoring.",
            "The scope excludes authentication and RBAC in the current phase, real external communication credentials, legal transplant workflow enforcement, and model training pipelines.",
            "The system is suitable for academic demonstration, functional validation, and architecture-level review for real-world expansion.",
        ],
        [
            "In-scope: workflow modules, APIs, dashboard metrics, database persistence",
            "Out-of-scope: auth, legal compliance automation, production cloud hardening",
            "Target usage: coordinator operations and executive monitoring",
        ],
    ),
    Chapter(
        "7. Literature Review",
        [
            "Prior studies in healthcare informatics emphasize interoperability and workflow orchestration as key factors in reducing critical care delays. Integrated platforms consistently outperform isolated systems in continuity and responsiveness.",
            "Research on donor matching highlights the importance of explainable criteria, particularly blood compatibility, geography, and availability windows, before introducing advanced machine learning ranking.",
            "Modern SaaS architecture literature recommends modular APIs, strong schema governance, and event logging for safety-critical coordination systems. These principles directly informed the LifeLink AI design.",
        ],
        [
            "Integrated workflow systems improve coordination quality",
            "Explainability is essential for medical decision confidence",
            "Relational schema integrity is foundational for healthcare applications",
            "Modular services simplify regulated system evolution",
        ],
    ),
    Chapter(
        "8. System Architecture",
        [
            "The architecture follows a layered model. The frontend (Next.js + TypeScript) consumes REST APIs exposed by NestJS backend modules. Backend services validate DTOs, execute business rules, and persist data via Prisma ORM into PostgreSQL.",
            "Cross-cutting concerns include Swagger-based API discoverability, strict relational constraints in schema design, and Dockerized runtime for repeatable setup.",
            "The architecture separates core operational logic from mock adapters so that messaging and AI extraction providers can be replaced with production integrations later.",
        ],
        [
            "Presentation Layer: Next.js App Router, React Query, Zustand",
            "Application Layer: NestJS modules with DTO validation",
            "Data Layer: PostgreSQL with Prisma relations and indexes",
            "Deployment Layer: Docker and docker-compose",
        ],
    ),
    Chapter(
        "9. Module Description",
        [
            "Patient Management module captures demographic profile, medical summary, guardian details, request urgency, and timeline placeholders.",
            "Donor Management module stores donor profile, health constraints, availability, and donation preferences. It supports eligibility-aware matching workflows.",
            "Matching module provides score breakdown, review notes, status progression, and relation tracking between patient and donor records. Notifications module logs dispatch events while report and blood bank modules strengthen operational readiness.",
        ],
        [
            "Patient module: list, detail, create, edit, request context",
            "Donor module: list, detail, create, edit, availability profiles",
            "Match module: score, status, review actions, traceability",
            "Notifications: in-app + mocked SMS/Email/WhatsApp with logs",
            "Reports: metadata, extraction placeholders, summary display",
            "Blood Bank: inventory, requests, low-stock alerts",
        ],
    ),
    Chapter(
        "10. Database Design",
        [
            "The database uses PostgreSQL with normalized relational design. Core entities include hospitals, coordinators, patients, donors, matches, report files, extractions, blood inventory, and notification logs.",
            "Foreign keys enforce referential consistency across patient requests, donor medical profiles, match reviews, and blood requests. Indexes are applied on search and workflow columns to optimize list and dashboard queries.",
            "Timestamp fields (`created_at`, `updated_at`) provide full audit context for operational actions and review steps.",
        ],
        [
            "Core tables: hospitals, coordinators, patients, donors, donor_patient_matches",
            "Clinical support: patient_medical_profiles, donor_medical_profiles, report_files, report_extractions",
            "Operations: notifications, notification_logs, blood_inventory, blood_requests, case_timelines",
            "Governance: strict PK/FK rules and indexed workflow columns",
        ],
    ),
    Chapter(
        "11. Matching Algorithm Explanation",
        [
            "The project implements a rule-based scoring engine as an AI-ready foundation. A donor-patient pair is evaluated across blood compatibility, donation type alignment, location proximity, donor availability, and medical risk indicators.",
            "Each factor contributes to component scores: `blood_compatibility_score`, `location_score`, and `availability_score`. The `overall_score` is computed through weighted aggregation and stored with every match record.",
            "This approach enables explainability. Coordinators can review why a score is high or low and override outcomes using structured review actions when clinical judgment requires intervention.",
        ],
        [
            "Blood group compatibility has highest weight",
            "Same district and city increases location score",
            "Unavailable donor status reduces availability contribution",
            "Medical flags apply score penalties",
            "Final score supports shortlist/approve/reject decisions",
        ],
    ),
    Chapter(
        "12. Technology Stack",
        [
            "Frontend is built with Next.js App Router, TypeScript, Tailwind CSS, and reusable UI components. React Query handles API state while Zustand supports client-level UI state.",
            "Backend uses NestJS with modular services, DTO validation, and Swagger documentation. Prisma provides typed database access on PostgreSQL.",
            "Docker and docker-compose provide reproducible local and demo environments for frontend, backend, and database services.",
        ],
        [
            "Frontend: Next.js, TypeScript, Tailwind CSS, React Query, Zustand",
            "Backend: NestJS, REST APIs, class-validator DTOs, Swagger",
            "Database: PostgreSQL, Prisma ORM",
            "DevOps: Docker, docker-compose",
        ],
    ),
    Chapter(
        "13. Implementation",
        [
            "Implementation followed phased delivery. Step 1 established foundational modules and schema. Step 2 added matching workflow upgrades, notifications abstraction, reports metadata and extraction placeholders, blood bank management, and dashboard integration.",
            "Frontend implementation emphasizes reusable cards, tables, badges, and form sections to maintain consistency across modules while reducing maintenance overhead.",
            "Backend implementation preserves backward compatibility by extending service layers and schema entities rather than breaking existing API contracts.",
        ],
        [
            "Schema-first incremental development",
            "Module-wise API implementation and integration",
            "Seeded demo data for realistic end-to-end walkthrough",
            "Consistent UI component reuse across pages",
        ],
    ),
    Chapter(
        "14. Testing and Results",
        [
            "Testing covered startup validation, API correctness, persistence checks, match score behavior, notification dispatch logging, report extraction placeholders, blood inventory alerts, and dashboard metric alignment.",
            "Functional API tests confirmed successful CRUD operations for patients and donors, match lifecycle updates, report creation and extraction response, and blood inventory access flows.",
            "Observed results matched expected outcomes for the implemented workflow scope. Validation failures returned explicit DTO errors, and persisted data was verified in PostgreSQL.",
        ],
        [
            "API routes respond with expected status codes and payloads",
            "Score breakdown fields persist with match records",
            "Notification logs recorded for triggered events",
            "Dashboard reflects underlying module counts and alerts",
        ],
    ),
    Chapter(
        "15. Advantages",
        [
            "LifeLink AI delivers a unified operational console that reduces coordination friction between donor, patient, and blood bank teams.",
            "Rule-based scoring improves decision clarity by making compatibility rationale visible to coordinators and executives.",
            "The modular architecture enables future feature growth without destabilizing current workflows.",
        ],
        [
            "Single source of operational truth",
            "Faster case review and shortlist cycles",
            "Explainable matching foundation",
            "Extendable API and schema design",
        ],
    ),
    Chapter(
        "16. Limitations",
        [
            "Current phase intentionally excludes authentication and production-grade security hardening. Access control, audit compliance, and user tenancy policies are future additions.",
            "Notification channels are mocked and do not use live credentials. Report extraction is placeholder logic and does not process actual report files.",
            "Geographic scoring is district-level and can be improved with route-aware travel estimation.",
        ],
        [
            "No auth/RBAC in current build",
            "Mock adapters for external communication",
            "AI extraction placeholders only",
            "No legal transplant workflow automation yet",
        ],
    ),
    Chapter(
        "17. Future Scope",
        [
            "Future roadmap includes authentication with RBAC, policy-driven approvals, production communication integrations, and model-assisted ranking using clinical and operational signals.",
            "A map-based donor finder and advanced analytics panel can improve region-level planning and emergency response.",
            "Compliance and legal workflow features can be introduced with configurable rulesets, digital consent trails, and audit exports.",
        ],
        [
            "Auth and multi-role permission framework",
            "AI-assisted matching and report interpretation",
            "Map visualization and geospatial ranking",
            "Forecasting for blood stock demand",
            "Regulatory and legal workflow modules",
        ],
    ),
    Chapter(
        "18. Conclusion",
        [
            "LifeLink AI demonstrates a complete full-stack coordination platform that addresses a real healthcare operations gap. The project transitions donor-patient matching from fragmented manual handling to a structured, reviewable workflow.",
            "The implemented system is academically complete for final-year evaluation and technically mature enough for controlled pilot demonstrations. Core workflows are functional, database-backed, and built using modern SaaS engineering practices.",
            "The architecture provides a strong base for advanced intelligence and production deployment phases.",
        ],
        [
            "Problem addressed with practical full-stack solution",
            "Workflow transparency improved through scores and logs",
            "System ready for next-phase productization",
        ],
    ),
    Chapter(
        "19. References",
        [
            "This project references standard software engineering, healthcare informatics, and modern web architecture practices. Sources are listed in APA-like format for academic review.",
        ],
        [
            "Pressman, R. S. Software Engineering: A Practitioner's Approach.",
            "Sommerville, I. Software Engineering.",
            "Fielding, R. Architectural Styles and the Design of Network-based Software Architectures.",
            "NIST Digital Health Guidelines and Data Governance Notes.",
            "Next.js Official Documentation.",
            "NestJS Official Documentation.",
            "Prisma ORM Documentation.",
            "PostgreSQL Official Documentation.",
            "Docker Documentation.",
            "Selected healthcare workflow and interoperability research papers.",
        ],
    ),
    Chapter(
        "20. Appendix",
        [
            "Appendix includes the user manual, testing evidence summary, sample API route list, and viva question bank used for oral defense preparation.",
        ],
        [
            "Appendix A: User Manual",
            "Appendix B: Testing Report and Result Matrix",
            "Appendix C: Viva Questions and Suggested Answers",
            "Appendix D: UI Screenshot Placeholders for final printed submission",
        ],
    ),
]

USER_MANUAL_POINTS = [
    "Prerequisites: Node.js, npm, Docker, docker-compose, and PostgreSQL access through containerized setup.",
    "Start with docker-compose up --build to run frontend, backend, and PostgreSQL services.",
    "Open the app at http://localhost:3000 and API docs at http://localhost:3001/api/docs.",
    "Use Patients module to add a patient profile with medical details, request details, and guardian information.",
    "Use Donors module to register donor profile, availability, and donation preferences.",
    "Use Matches module to review calculated compatibility score and perform review actions.",
    "Use Reports module to upload report metadata and trigger mock extraction summary.",
    "Use Notifications module to monitor read/unread status and dispatch logs.",
    "Use Blood Banks module to check group-wise stock and low stock indicators.",
    "Use Dashboard for executive-level totals, pending reviews, alerts, and activity feed.",
]

VIVA_QA = [
    ("What is the main aim of LifeLink AI?", "To provide an integrated platform for donor-patient coordination with explainable matching workflows and operational visibility."),
    ("Why did you choose NestJS for backend?", "NestJS provides modular architecture, DTO validation, and clear controller-service separation suited for scalable APIs."),
    ("How is compatibility score calculated?", "Using weighted rule-based factors including blood compatibility, location proximity, donor availability, and health risk penalties."),
    ("Why is Prisma used in this project?", "Prisma provides typed ORM access, migration support, and maintainable schema-driven development for PostgreSQL."),
    ("How do you ensure API correctness?", "By using DTO validation, Swagger contracts, and endpoint-level functional testing with expected versus actual comparisons."),
    ("What is the role of React Query?", "React Query manages server-state fetching, caching, invalidation, and loading/error states in the frontend."),
    ("How are notifications implemented?", "Through in-app storage with channel abstractions (SMS/Email/WhatsApp) currently mocked for demo-safe operation."),
    ("What are major limitations in this phase?", "No auth/RBAC, no live communication credentials, and no production-grade AI extraction pipeline."),
    ("How is future AI integration planned?", "By keeping matching and extraction logic modular so placeholder services can be replaced by model-backed services."),
    ("How did you validate persistence?", "By checking API results and confirming records in PostgreSQL tables for patients, donors, matches, and logs."),
    ("What deployment approach is used?", "Dockerized services with docker-compose for reproducible local execution."),
    ("How is dashboard data generated?", "From aggregated backend queries across patient, donor, match, report, notification, and blood inventory modules."),
]

TEST_CASES = [
    ("TC-01", "Create patient with minimum valid fields", "POST /patients", "201 Created and persisted record", "Pass"),
    ("TC-02", "Create patient with detailed profile", "POST /patients", "201 Created with nested records", "Pass"),
    ("TC-03", "Create patient with invalid hospital UUID", "POST /patients", "400 with validation message", "Pass"),
    ("TC-04", "Create donor with availability and preferences", "POST /donors", "201 Created and persisted record", "Pass"),
    ("TC-05", "Fetch dashboard summary", "GET /dashboard/summary", "200 with aggregate metrics", "Pass"),
    ("TC-06", "Create match and score breakdown", "POST /matches", "201 with score fields", "Pass"),
    ("TC-07", "Update match status", "PATCH /matches/{id}/status", "200 and status updated", "Pass"),
    ("TC-08", "Add review action", "POST /matches/{id}/reviews", "200 and review appended", "Pass"),
    ("TC-09", "Create notification record", "POST /notifications", "201 and log entry generated", "Pass"),
    ("TC-10", "Upload report metadata", "POST /reports", "201 and report listed", "Pass"),
    ("TC-11", "Trigger mock extraction", "POST /reports/{id}/extract", "200 with extraction summary", "Pass"),
    ("TC-12", "Get low stock alerts", "GET /blood-banks/low-stock-alerts", "200 with alert rows", "Pass"),
]

SLIDES = [
    ("LifeLink AI", [
        "Intelligent Donor-Patient Matching and Healthcare Coordination System",
        "Final Year Project Presentation",
        "Department of Computer Science and Engineering",
    ]),
    ("Problem Context", [
        "Hospitals manage donor and patient workflows across disconnected tools",
        "Critical case delays occur due to manual matching and communication",
        "Limited visibility of blood stock risk and review bottlenecks",
    ]),
    ("Project Vision", [
        "Create one integrated healthcare coordination platform",
        "Enable explainable matching with operational dashboards",
        "Design an AI-ready architecture for future intelligent upgrades",
    ]),
    ("Objectives", [
        "Centralize patient, donor, and match lifecycle data",
        "Improve decision speed with score breakdowns",
        "Provide end-to-end traceability and notifications",
    ]),
    ("System Architecture", [
        "Frontend: Next.js + TypeScript + Tailwind CSS",
        "Backend: NestJS REST APIs with DTO validation",
        "Database: PostgreSQL with Prisma ORM",
        "Deployment: Docker and docker-compose",
    ]),
    ("Core Modules", [
        "Patient management",
        "Donor management",
        "Match workflow and review",
        "Notifications, reports, blood bank, dashboard",
    ]),
    ("Database Design", [
        "Normalized relational schema with PK/FK constraints",
        "Index-driven query optimization",
        "Audit fields for lifecycle and compliance readiness",
    ]),
    ("Matching Logic", [
        "Rule-based compatibility engine",
        "Factors: blood group, location, availability, donation type, risk flags",
        "Stored score breakdown supports explainable review",
    ]),
    ("Patient and Donor Flows", [
        "Structured forms with validation",
        "Medical profile and request details captured",
        "Coordinator-ready list and detail views",
    ]),
    ("Match Review Experience", [
        "Actions: review, shortlist, approve, reject, notify",
        "Status transitions and review notes are logged",
        "Timeline placeholder for complete case traceability",
    ]),
    ("Reports and Intelligence Placeholder", [
        "Report metadata upload and listing",
        "Mock extraction service returns parsed summary fields",
        "UI shows extraction status and summarized insights",
    ]),
    ("Blood Bank Module", [
        "Group-wise stock visibility",
        "Low stock alerts and recent stock activity",
        "Basic blood request tracking",
    ]),
    ("Notification Abstraction", [
        "In-app notifications and dispatch logs",
        "Mock channel adapters: SMS, Email, WhatsApp",
        "Template-driven event notifications",
    ]),
    ("Dashboard Highlights", [
        "Counts: patients, donors, matches, pending reviews",
        "Alerts: low stock and urgent workflows",
        "Recent notifications and activity feed",
    ]),
    ("Testing Strategy", [
        "Startup, API, database, workflow, and dashboard validation",
        "Expected versus actual test matrix maintained",
        "Regression checks for backward compatibility",
    ]),
    ("Results", [
        "All primary modules integrated end-to-end",
        "Stable API behavior with validation errors surfaced correctly",
        "Persistent records confirmed in PostgreSQL",
    ]),
    ("Limitations", [
        "Authentication and RBAC intentionally deferred",
        "External communication adapters are mocked",
        "AI extraction currently placeholder-driven",
    ]),
    ("Future Scope", [
        "Auth and role governance",
        "Advanced AI matching and extraction",
        "Map-based donor finder and predictive analytics",
    ]),
    ("Conclusion", [
        "LifeLink AI delivers a robust healthcare coordination foundation",
        "Architecture is scalable for product-level evolution",
        "Demonstrates strong academic and practical relevance",
    ]),
    ("Thank You", [
        "Questions and viva discussion",
        "Prepared by: Project Team",
    ]),
]


def set_doc_defaults(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    style.font.size = Pt(12)


def add_centered_paragraph(doc: Document, text: str, size: int = 12, bold: bool = False) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Times New Roman"


def add_docx_title_pages(doc: Document) -> None:
    add_centered_paragraph(doc, "FINAL YEAR PROJECT REPORT", size=16, bold=True)
    doc.add_paragraph()
    add_centered_paragraph(doc, PROJECT_TITLE, size=16, bold=True)
    doc.add_paragraph()
    add_centered_paragraph(doc, "Submitted in partial fulfillment of the requirements", size=12)
    add_centered_paragraph(doc, "for the award of the degree of Bachelor of Engineering", size=12)
    add_centered_paragraph(doc, "in Computer Science and Engineering", size=12)
    doc.add_paragraph()
    add_centered_paragraph(doc, "Submitted By", size=12, bold=True)
    add_centered_paragraph(doc, "Student Name: _________________________", size=12)
    add_centered_paragraph(doc, "Register Number: ______________________", size=12)
    doc.add_paragraph()
    add_centered_paragraph(doc, "Under the Guidance of", size=12, bold=True)
    add_centered_paragraph(doc, "Guide Name: ___________________________", size=12)
    doc.add_paragraph()
    add_centered_paragraph(doc, "Department of Computer Science and Engineering", size=12)
    add_centered_paragraph(doc, "College Name: _________________________", size=12)
    add_centered_paragraph(doc, f"Date: {TODAY}", size=12)

    doc.add_page_break()
    doc.add_heading("Bonafide Certificate", level=1)
    doc.add_paragraph(
        "This is to certify that the project report entitled \"LifeLink AI - Intelligent Donor-Patient "
        "Matching and Healthcare Coordination System\" is a bonafide work carried out by the student "
        "of final year B.E. Computer Science and Engineering under my guidance and supervision."
    )
    doc.add_paragraph("\nGuide Signature: ________________________")
    doc.add_paragraph("HOD Signature: __________________________")
    doc.add_paragraph("Principal Signature: _____________________")

    doc.add_page_break()
    doc.add_heading("Declaration", level=1)
    doc.add_paragraph(
        "I hereby declare that this project report is the original work carried out by me and has not "
        "been submitted in part or full for any other degree or diploma in this or any other institution."
    )
    doc.add_paragraph("\nStudent Signature: ______________________")

    doc.add_page_break()
    doc.add_heading("Acknowledgement", level=1)
    doc.add_paragraph(
        "I express my sincere gratitude to the management, Head of Department, project guide, faculty "
        "members, and peers for their continuous encouragement and valuable guidance throughout the "
        "development of this project. I also thank my family for their support and motivation."
    )

    doc.add_page_break()
    doc.add_heading("Abstract", level=1)
    doc.add_paragraph(
        "LifeLink AI is a full-stack healthcare coordination platform built to streamline donor-patient "
        "matching, hospital case review, report monitoring, notifications, and blood bank operations. "
        "The system integrates patient and donor workflows with a rule-based compatibility engine that "
        "generates explainable score breakdowns. A NestJS backend with PostgreSQL and Prisma ensures "
        "strong data consistency, while a Next.js frontend provides responsive workflow interfaces and "
        "executive dashboards. The implementation demonstrates a scalable architecture where AI modules "
        "can be incorporated in future phases without disrupting operational reliability."
    )


def add_docx_toc_and_lists(doc: Document) -> None:
    doc.add_page_break()
    doc.add_heading("Table of Contents", level=1)
    toc_items = [
        "Title Page",
        "Bonafide Certificate",
        "Declaration",
        "Acknowledgement",
        "Abstract",
        "List of Figures",
        "List of Tables",
    ] + [chapter.title for chapter in CHAPTERS] + [
        "Appendix A - User Manual",
        "Appendix B - Testing Report",
        "Appendix C - Viva Questions and Answers",
    ]
    for item in toc_items:
        doc.add_paragraph(item, style="List Number")

    doc.add_page_break()
    doc.add_heading("List of Figures", level=1)
    figures = [
        "Figure 1: High-level system architecture",
        "Figure 2: Dashboard and executive metrics view",
        "Figure 3: Match workflow state transition model",
        "Figure 4: Database relationship overview",
        "Figure 5: Blood inventory and low stock alert panel",
        "Figure 6: Report extraction summary view",
    ]
    for fig in figures:
        doc.add_paragraph(fig, style="List Number")

    doc.add_page_break()
    doc.add_heading("List of Tables", level=1)
    tables = [
        "Table 1: Technology stack summary",
        "Table 2: Module responsibility matrix",
        "Table 3: Core database tables",
        "Table 4: API endpoint summary",
        "Table 5: Test case matrix",
        "Table 6: Viva question bank",
    ]
    for tab in tables:
        doc.add_paragraph(tab, style="List Number")


def add_docx_chapters(doc: Document) -> None:
    for chapter in CHAPTERS:
        doc.add_page_break()
        doc.add_heading(chapter.title, level=1)
        for para in chapter.paragraphs:
            doc.add_paragraph(para)
        if chapter.bullets:
            for bullet in chapter.bullets:
                doc.add_paragraph(bullet, style="List Bullet")

        if chapter.title.startswith("12."):
            table = doc.add_table(rows=1, cols=2)
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            hdr[0].text = "Layer"
            hdr[1].text = "Technologies"
            rows = [
                ("Frontend", "Next.js, TypeScript, Tailwind CSS, React Query, Zustand"),
                ("Backend", "NestJS, REST API, DTO validation, Swagger"),
                ("Database", "PostgreSQL, Prisma ORM"),
                ("Deployment", "Docker, docker-compose"),
            ]
            for left, right in rows:
                r = table.add_row().cells
                r[0].text = left
                r[1].text = right

        if chapter.title.startswith("10."):
            table = doc.add_table(rows=1, cols=3)
            table.style = "Table Grid"
            table.rows[0].cells[0].text = "Table Name"
            table.rows[0].cells[1].text = "Purpose"
            table.rows[0].cells[2].text = "Key Fields"
            db_rows = [
                ("patients", "Stores patient profile", "hospital_id, uhid, urgency_level, case_status"),
                ("donors", "Stores donor profile", "hospital_id, donor_code, status"),
                ("donor_patient_matches", "Stores match scores", "patient_id, donor_id, overall_score, status"),
                ("report_files", "Stores report metadata", "patient_id, file_type, extraction_status"),
                ("notification_logs", "Stores dispatch trace", "channel, status, event_type"),
                ("blood_inventory", "Stores stock units", "blood_bank_id, blood_group, units_available"),
            ]
            for row in db_rows:
                c = table.add_row().cells
                c[0].text, c[1].text, c[2].text = row

        if chapter.title.startswith("14."):
            doc.add_heading("Test Plan", level=2)
            doc.add_paragraph(
                "Testing was executed in staged order: startup verification, API validation, data persistence "
                "checks, matching workflow validation, notification trigger validation, report extraction checks, "
                "blood bank alert checks, and dashboard verification."
            )
            doc.add_heading("Test Case Matrix", level=2)
            table = doc.add_table(rows=1, cols=5)
            table.style = "Table Grid"
            headers = ["Test ID", "Scenario", "Endpoint/Area", "Expected Result", "Status"]
            for idx, text in enumerate(headers):
                table.rows[0].cells[idx].text = text
            for row in TEST_CASES:
                cells = table.add_row().cells
                for i, value in enumerate(row):
                    cells[i].text = value


def add_docx_appendix(doc: Document) -> None:
    doc.add_page_break()
    doc.add_heading("Appendix A - User Manual", level=1)
    for item in USER_MANUAL_POINTS:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_page_break()
    doc.add_heading("Appendix B - Testing Report Summary", level=1)
    doc.add_paragraph(
        "The system was tested for functional correctness, persistence reliability, score behavior, "
        "workflow transitions, and dashboard consistency. All primary workflows passed in the test environment."
    )
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Validation Area"
    table.rows[0].cells[1].text = "Method"
    table.rows[0].cells[2].text = "Outcome"
    test_summary_rows = [
        ("API validation", "Swagger and curl endpoint checks", "All core endpoints returned expected status"),
        ("Database persistence", "Record verification in PostgreSQL", "Inserted records persisted correctly"),
        ("Matching logic", "Rule-based scenario checks", "Score pattern aligned with expected criteria"),
        ("Notifications", "Trigger and log checks", "Notification logs recorded per dispatch"),
        ("Dashboard", "Aggregate metric comparison", "Counts and alerts matched backend data"),
    ]
    for row in test_summary_rows:
        c = table.add_row().cells
        c[0].text, c[1].text, c[2].text = row

    doc.add_page_break()
    doc.add_heading("Appendix C - Viva Questions and Answers", level=1)
    for q, a in VIVA_QA:
        p = doc.add_paragraph()
        p.add_run("Q: ").bold = True
        p.add_run(q)
        p2 = doc.add_paragraph()
        p2.add_run("A: ").bold = True
        p2.add_run(a)

    doc.add_page_break()
    doc.add_heading("Appendix D - UI Screenshot Placeholders", level=1)
    placeholders = [
        "Screenshot Placeholder 1: Executive dashboard with KPI cards",
        "Screenshot Placeholder 2: Patient create and detail workflow",
        "Screenshot Placeholder 3: Donor registry and profile detail",
        "Screenshot Placeholder 4: Match list and score breakdown view",
        "Screenshot Placeholder 5: Reports list and extraction summary",
        "Screenshot Placeholder 6: Blood inventory and low stock alerts",
    ]
    for text in placeholders:
        doc.add_paragraph(text, style="List Number")


def create_documentation_docx(path: Path) -> None:
    doc = Document()
    set_doc_defaults(doc)
    add_docx_title_pages(doc)
    add_docx_toc_and_lists(doc)
    add_docx_chapters(doc)
    add_docx_appendix(doc)
    doc.save(path)


def _pdf_styles():
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "CustomTitle",
            parent=base["Title"],
            alignment=TA_CENTER,
            fontName="Times-Bold",
            fontSize=20,
            leading=24,
            spaceAfter=14,
        ),
        "h1": ParagraphStyle(
            "CustomH1",
            parent=base["Heading1"],
            fontName="Times-Bold",
            fontSize=15,
            leading=19,
            spaceBefore=10,
            spaceAfter=8,
        ),
        "h2": ParagraphStyle(
            "CustomH2",
            parent=base["Heading2"],
            fontName="Times-Bold",
            fontSize=12,
            leading=15,
            spaceBefore=8,
            spaceAfter=6,
        ),
        "body": ParagraphStyle(
            "CustomBody",
            parent=base["BodyText"],
            fontName="Times-Roman",
            fontSize=11,
            leading=15,
            alignment=TA_JUSTIFY,
            spaceAfter=6,
        ),
        "left": ParagraphStyle(
            "CustomLeft",
            parent=base["BodyText"],
            fontName="Times-Roman",
            fontSize=11,
            leading=14,
            alignment=TA_LEFT,
            spaceAfter=5,
        ),
    }


def build_documentation_pdf(path: Path) -> None:
    styles = _pdf_styles()
    story = []

    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph("FINAL YEAR PROJECT REPORT", styles["title"]))
    story.append(Paragraph(PROJECT_TITLE, styles["title"]))
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Department of Computer Science and Engineering", styles["left"]))
    story.append(Paragraph(f"Submission Date: {TODAY}", styles["left"]))
    story.append(PageBreak())

    front_sections = [
        (
            "Bonafide Certificate",
            "This is to certify that this project work has been carried out by final-year students under authorized guidance and supervision for academic submission.",
        ),
        (
            "Declaration",
            "The project report is declared as original work and has not been submitted elsewhere for another degree or diploma.",
        ),
        (
            "Acknowledgement",
            "The project team gratefully acknowledges the support of faculty, guide, peers, and family during the design and implementation process.",
        ),
        (
            "Abstract",
            "LifeLink AI is an integrated healthcare coordination platform combining patient management, donor management, matching workflows, report processing placeholders, notification abstraction, blood bank visibility, and executive dashboards in one scalable system.",
        ),
    ]
    for heading, text in front_sections:
        story.append(Paragraph(heading, styles["h1"]))
        story.append(Paragraph(text, styles["body"]))
        story.append(PageBreak())

    story.append(Paragraph("Table of Contents", styles["h1"]))
    for idx, chapter in enumerate(CHAPTERS, start=1):
        story.append(Paragraph(f"{idx}. {chapter.title.split('. ', 1)[1]}", styles["left"]))
    story.append(PageBreak())

    story.append(Paragraph("List of Figures", styles["h1"]))
    for fig in [
        "Figure 1: High-level system architecture",
        "Figure 2: Dashboard module overview",
        "Figure 3: Match workflow and review lifecycle",
        "Figure 4: Database ER relation summary",
        "Figure 5: Blood inventory low stock view",
    ]:
        story.append(Paragraph(fig, styles["left"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph("List of Tables", styles["h1"]))
    for tab in [
        "Table 1: Technology stack",
        "Table 2: Core module mapping",
        "Table 3: Database table summary",
        "Table 4: API and workflow test matrix",
    ]:
        story.append(Paragraph(tab, styles["left"]))
    story.append(PageBreak())

    for chapter in CHAPTERS:
        story.append(Paragraph(chapter.title, styles["h1"]))
        for p in chapter.paragraphs:
            story.append(Paragraph(p, styles["body"]))
        for b in chapter.bullets:
            story.append(Paragraph(f"- {b}", styles["left"]))

        if chapter.title.startswith("12."):
            data = [
                ["Layer", "Technologies"],
                ["Frontend", "Next.js, TypeScript, Tailwind CSS, React Query, Zustand"],
                ["Backend", "NestJS, REST API, DTO validation, Swagger"],
                ["Database", "PostgreSQL, Prisma ORM"],
                ["Deployment", "Docker, docker-compose"],
            ]
            t = Table(data, colWidths=[1.8 * inch, 4.9 * inch])
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Times-Roman"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
            ]))
            story.append(Spacer(1, 8))
            story.append(t)

        if chapter.title.startswith("14."):
            story.append(Paragraph("Test Case Matrix", styles["h2"]))
            data = [["Test ID", "Scenario", "Area", "Expected", "Result"]] + [list(tc) for tc in TEST_CASES]
            t = Table(data, colWidths=[0.7 * inch, 1.8 * inch, 1.5 * inch, 2.3 * inch, 0.8 * inch])
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.4, colors.black),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Times-Roman"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(Spacer(1, 8))
            story.append(t)

        story.append(PageBreak())

    story.append(Paragraph("Appendix A - User Manual", styles["h1"]))
    for item in USER_MANUAL_POINTS:
        story.append(Paragraph(f"- {item}", styles["left"]))
    story.append(PageBreak())

    story.append(Paragraph("Appendix C - Viva Questions and Answers", styles["h1"]))
    for q, a in VIVA_QA:
        story.append(Paragraph(f"Q: {q}", styles["h2"]))
        story.append(Paragraph(f"A: {a}", styles["body"]))

    def on_page(canv: canvas.Canvas, doc_obj):
        canv.setFont("Times-Roman", 9)
        canv.drawRightString(A4[0] - 40, 20, f"Page {canv.getPageNumber()}")

    pdf = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=40,
        rightMargin=40,
        topMargin=40,
        bottomMargin=30,
    )
    pdf.build(story, onFirstPage=on_page, onLaterPages=on_page)


def create_synopsis_docx(path: Path) -> None:
    doc = Document()
    set_doc_defaults(doc)

    add_centered_paragraph(doc, "PROJECT SYNOPSIS", size=16, bold=True)
    add_centered_paragraph(doc, PROJECT_TITLE, size=14, bold=True)
    add_centered_paragraph(doc, f"Date: {TODAY}", size=12)

    doc.add_heading("1. Introduction", level=1)
    doc.add_paragraph(
        "LifeLink AI is a full-stack healthcare coordination platform developed to integrate patient, donor, "
        "matching, reporting, notifications, and blood bank operations into one unified system."
    )

    doc.add_heading("2. Problem Statement", level=1)
    doc.add_paragraph(
        "Current hospital coordination workflows are fragmented and mostly manual, causing delays in donor-patient "
        "matching and reducing operational transparency."
    )

    doc.add_heading("3. Proposed Solution", level=1)
    for b in [
        "Unified dashboard for coordinator and executive views",
        "Rule-based donor-patient compatibility scoring",
        "Match review lifecycle with notes and status changes",
        "Report metadata and extraction placeholder workflows",
        "Blood inventory and low-stock alert visibility",
    ]:
        doc.add_paragraph(b, style="List Bullet")

    doc.add_heading("4. Objectives", level=1)
    for b in [
        "Reduce time to shortlist suitable donors",
        "Improve workflow traceability and communication",
        "Provide scalable architecture for future AI integration",
    ]:
        doc.add_paragraph(b, style="List Bullet")

    doc.add_heading("5. Technology Stack", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Layer"
    table.rows[0].cells[1].text = "Tools"
    rows = [
        ("Frontend", "Next.js, TypeScript, Tailwind CSS"),
        ("Backend", "NestJS, REST APIs, Swagger"),
        ("Database", "PostgreSQL, Prisma"),
        ("Deployment", "Docker, docker-compose"),
    ]
    for r in rows:
        c = table.add_row().cells
        c[0].text, c[1].text = r

    doc.add_heading("6. Expected Outcome", level=1)
    doc.add_paragraph(
        "A production-style healthcare SaaS foundation that demonstrates complete donor-patient coordination "
        "workflows with reliable persistence, testing readiness, and extensibility for future intelligent services."
    )

    doc.add_heading("7. Conclusion", level=1)
    doc.add_paragraph(
        "The project delivers a practical and academically relevant system that addresses real healthcare workflow "
        "challenges with modern full-stack engineering practices."
    )

    doc.save(path)


def create_synopsis_pdf(path: Path) -> None:
    styles = _pdf_styles()
    story = [
        Paragraph("PROJECT SYNOPSIS", styles["title"]),
        Paragraph(PROJECT_TITLE, styles["h1"]),
        Paragraph(f"Date: {TODAY}", styles["left"]),
        Spacer(1, 12),
        Paragraph("Introduction", styles["h1"]),
        Paragraph(
            "LifeLink AI integrates donor-patient matching and hospital coordination workflows in a modular "
            "full-stack web platform.",
            styles["body"],
        ),
        Paragraph("Problem Statement", styles["h1"]),
        Paragraph(
            "Manual and disconnected systems delay critical matching decisions and reduce coordination quality.",
            styles["body"],
        ),
        Paragraph("Proposed Solution", styles["h1"]),
    ]
    for item in [
        "Unified workflow dashboard",
        "Rule-based compatibility score",
        "Notifications and report placeholders",
        "Blood inventory alerting",
    ]:
        story.append(Paragraph(f"- {item}", styles["left"]))

    story.extend([
        Paragraph("Technology Stack", styles["h1"]),
        Paragraph("Frontend: Next.js + TypeScript + Tailwind CSS", styles["left"]),
        Paragraph("Backend: NestJS + REST + Swagger", styles["left"]),
        Paragraph("Database: PostgreSQL + Prisma", styles["left"]),
        Paragraph("Deployment: Docker + docker-compose", styles["left"]),
        Paragraph("Expected Outcome", styles["h1"]),
        Paragraph(
            "A scalable healthcare coordination solution with reliable APIs, modular services, and practical "
            "readiness for future AI-assisted workflows.",
            styles["body"],
        ),
    ])

    pdf = SimpleDocTemplate(str(path), pagesize=A4, leftMargin=40, rightMargin=40, topMargin=40, bottomMargin=35)
    pdf.build(story)


def create_presentation_pptx(path: Path) -> None:
    prs = Presentation()

    # Theme colors
    title_color = RGBColor(16, 70, 135)
    body_color = RGBColor(30, 41, 59)

    for idx, (title, bullets) in enumerate(SLIDES):
        if idx == 0:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = "\n".join(bullets)
            for p in slide.placeholders[1].text_frame.paragraphs:
                p.font.size = PptPt(20 if p.level == 0 else 16)
                p.font.color.rgb = body_color
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
            continue

        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color

        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            p.level = 0
            p.font.size = PptPt(22 if len(bullets) <= 3 else 20)
            p.font.color.rgb = body_color

        # Add screenshot placeholder on module/visual slides
        if title in {
            "Patient and Donor Flows",
            "Match Review Experience",
            "Reports and Intelligence Placeholder",
            "Blood Bank Module",
            "Dashboard Highlights",
        }:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                PptInches(6.2),
                PptInches(1.6),
                PptInches(3.0),
                PptInches(3.8),
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(232, 241, 250)
            shape.line.color.rgb = RGBColor(126, 176, 222)
            text = shape.text_frame
            text.text = "UI Screenshot Placeholder"
            text.paragraphs[0].font.size = PptPt(16)
            text.paragraphs[0].font.bold = True
            text.paragraphs[0].font.color.rgb = RGBColor(40, 75, 120)
            p = text.add_paragraph()
            p.text = "Insert final captured screens here"
            p.font.size = PptPt(13)
            p.font.color.rgb = RGBColor(40, 75, 120)

    prs.save(path)


def create_presentation_pdf(path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=landscape(A4))
    width, height = landscape(A4)

    for title, bullets in SLIDES:
        c.setFillColorRGB(0.08, 0.31, 0.57)
        c.setFont("Helvetica-Bold", 28)
        c.drawString(45, height - 60, title)

        y = height - 110
        c.setFillColorRGB(0.12, 0.16, 0.24)
        c.setFont("Helvetica", 16)
        for item in bullets:
            c.drawString(60, y, f"- {item}")
            y -= 28

        if title in {
            "Patient and Donor Flows",
            "Match Review Experience",
            "Reports and Intelligence Placeholder",
            "Blood Bank Module",
            "Dashboard Highlights",
        }:
            c.setStrokeColorRGB(0.49, 0.68, 0.85)
            c.setFillColorRGB(0.92, 0.96, 0.99)
            c.roundRect(width - 300, 130, 250, 230, 12, fill=1, stroke=1)
            c.setFillColorRGB(0.2, 0.35, 0.52)
            c.setFont("Helvetica-Bold", 13)
            c.drawString(width - 278, 330, "Screenshot Placeholder")
            c.setFont("Helvetica", 11)
            c.drawString(width - 278, 308, "Insert module UI capture")

        c.setFont("Helvetica", 9)
        c.setFillColorRGB(0.35, 0.4, 0.45)
        c.drawRightString(width - 35, 18, "LifeLink AI - Final Presentation")
        c.showPage()

    c.save()


def create_one_page_summary_pdf(path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4

    c.setFont("Helvetica-Bold", 18)
    c.drawString(40, height - 55, "One Page Project Summary")

    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, height - 85, "LifeLink AI - Intelligent Donor-Patient Matching and Healthcare Coordination System")

    c.setFont("Helvetica", 11)
    y = height - 120
    lines = [
        "Project Type: Full-stack healthcare SaaS web application",
        "Core Modules: Patients, Donors, Matching, Notifications, Reports, Blood Bank, Dashboard",
        "Frontend Stack: Next.js, TypeScript, Tailwind CSS",
        "Backend Stack: NestJS with DTO validation and Swagger",
        "Database: PostgreSQL with Prisma ORM",
        "Deployment: Docker and docker-compose",
        "Matching Logic: Rule-based score using blood compatibility, location, availability, and risk factors",
        "Key Value: Unified operational workflow with explainable decision support",
        "Testing Coverage: API, persistence, score behavior, notifications, dashboard consistency",
        "Future Scope: Auth/RBAC, advanced AI scoring, map-based donor finder, production integrations",
    ]
    for line in lines:
        c.drawString(45, y, f"- {line}")
        y -= 23

    c.setFont("Helvetica-Bold", 12)
    c.drawString(40, 95, "Outcome")
    c.setFont("Helvetica", 11)
    c.drawString(45, 75, "A robust, modular, and academic-grade platform ready for evaluation and future expansion.")

    c.setFont("Helvetica", 9)
    c.drawRightString(width - 40, 20, f"Generated on {TODAY}")
    c.save()


def generate_all(output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)

    files = {
        "documentation_docx": output_dir / "LifeLink_AI_Final_Project_Documentation.docx",
        "documentation_pdf": output_dir / "LifeLink_AI_Final_Project_Documentation.pdf",
        "presentation_pptx": output_dir / "LifeLink_AI_Final_Presentation.pptx",
        "presentation_pdf": output_dir / "LifeLink_AI_Final_Presentation.pdf",
        "synopsis_docx": output_dir / "LifeLink_AI_Project_Synopsis.docx",
        "synopsis_pdf": output_dir / "LifeLink_AI_Project_Synopsis.pdf",
        "one_page_pdf": output_dir / "LifeLink_AI_One_Page_Summary.pdf",
    }

    create_documentation_docx(files["documentation_docx"])
    build_documentation_pdf(files["documentation_pdf"])
    create_presentation_pptx(files["presentation_pptx"])
    create_presentation_pdf(files["presentation_pdf"])
    create_synopsis_docx(files["synopsis_docx"])
    create_synopsis_pdf(files["synopsis_pdf"])
    create_one_page_summary_pdf(files["one_page_pdf"])

    return list(files.values())


def main() -> None:
    out = Path("/Users/rohith/AI-Doner/output/submission_pack")
    generated = generate_all(out)
    print("Generated files:")
    for f in generated:
        print(f"- {f}")


if __name__ == "__main__":
    main()
